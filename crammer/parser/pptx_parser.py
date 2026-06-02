from pptx import Presentation


def parse_pptx(file_path):
    try:
        prs = Presentation(file_path)
        pages = []
        for i, slide in enumerate(prs.slides):
            text_parts = []
            md_tables = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            text_parts.append(t)
                if shape.has_table:
                    tbl_data = []
                    for row in shape.table.rows:
                        tbl_data.append(
                            [cell.text for cell in row.cells]
                        )
                    md_tables.append(_table_to_md(tbl_data))
            pages.append(
                {
                    "page_num": i + 1,
                    "text": "\n".join(text_parts),
                    "tables": md_tables,
                }
            )
        return {"pages": pages, "total_pages": len(pages)}
    except Exception as e:
        return {"pages": [], "total_pages": 0, "error": str(e)}


def _table_to_md(table):
    if not table:
        return ""
    cleaned = [[str(c) if c is not None else "" for c in row] for row in table]
    col_count = max(len(row) for row in cleaned) if cleaned else 0
    rows = []
    for i, row in enumerate(cleaned):
        padded = row + [""] * (col_count - len(row))
        rows.append("| " + " | ".join(padded) + " |")
        if i == 0:
            rows.append("|" + "|".join(["---"] * col_count) + "|")
    return "\n".join(rows)
