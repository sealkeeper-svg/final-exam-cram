import pytest
from pptx import Presentation
from pptx.util import Inches
from crammer.parser.pptx_parser import parse_pptx


def _make_pptx(filepath, slides_data):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    for title, body, table_data in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        left = Inches(1)
        top = Inches(0.5)
        width = Inches(8)
        height = Inches(1)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = title
        body_box = slide.shapes.add_textbox(left, Inches(1.5), width, Inches(4))
        bf = body_box.text_frame
        bf.text = body
        if table_data:
            rows, cols = len(table_data), len(table_data[0]) if table_data else 0
            tbl_shape = slide.shapes.add_table(
                rows, cols, left, Inches(4), width, Inches(2)
            )
            for r in range(rows):
                for c in range(cols):
                    tbl_shape.table.cell(r, c).text = str(table_data[r][c])
    prs.save(str(filepath))


def test_parse_pptx_single_slide(tmp_path):
    path = tmp_path / "test.pptx"
    _make_pptx(str(path), [("标题", "正文内容", None)])
    result = parse_pptx(str(path))
    assert result["total_pages"] == 1
    page = result["pages"][0]
    assert page["page_num"] == 1
    assert "标题" in page["text"]
    assert "正文内容" in page["text"]
    assert page["tables"] == []


def test_parse_pptx_with_table(tmp_path):
    path = tmp_path / "test_table.pptx"
    _make_pptx(str(path), [("标题", "正文", [["A", "B"], ["1", "2"]])])
    result = parse_pptx(str(path))
    assert result["total_pages"] == 1
    page = result["pages"][0]
    assert len(page["tables"]) >= 1


def test_parse_pptx_multi_slide(tmp_path):
    path = tmp_path / "multi.pptx"
    _make_pptx(
        str(path),
        [("Slide 1", "Content 1", None), ("Slide 2", "Content 2", None)],
    )
    result = parse_pptx(str(path))
    assert result["total_pages"] == 2
    assert result["pages"][0]["page_num"] == 1
    assert result["pages"][1]["page_num"] == 2


def test_parse_pptx_missing_file():
    result = parse_pptx("/nonexistent/file_12345.pptx")
    assert result["total_pages"] == 0
    assert result["pages"] == []
    assert "error" in result
